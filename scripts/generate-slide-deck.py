"""
Generate PowerPoint deck from slide outline.
Output: decks/agent-fidelity-spectrum.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
import os

# ── Colors ──────────────────────────────────────────────────────────────
DARK_NAVY = RGBColor(0x1B, 0x2A, 0x4A)
MEDIUM_BLUE = RGBColor(0x2C, 0x5F, 0x8A)
ACCENT_BLUE = RGBColor(0x35, 0x7A, 0xBD)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT_GRAY = RGBColor(0xF2, 0xF2, 0xF2)
DARK_TEXT = RGBColor(0x2D, 0x2D, 0x2D)
MEDIUM_TEXT = RGBColor(0x55, 0x55, 0x55)
TABLE_HEADER_BG = RGBColor(0x2C, 0x5F, 0x8A)
TABLE_ALT_BG = RGBColor(0xE8, 0xF0, 0xF8)

# Level colors (green to red spectrum)
L1_GREEN = RGBColor(0x4C, 0xAF, 0x50)
L2_LIME = RGBColor(0x8B, 0xC3, 0x4A)
L3_ORANGE = RGBColor(0xFF, 0x98, 0x00)
L4_DEEP_ORANGE = RGBColor(0xF4, 0x51, 0x1E)
L5_RED = RGBColor(0xE5, 0x39, 0x35)
LEVEL_COLORS = [L1_GREEN, L2_LIME, L3_ORANGE, L4_DEEP_ORANGE, L5_RED]

# Severity colors
SEV_CRITICAL = RGBColor(0xC6, 0x28, 0x28)
SEV_HIGH = RGBColor(0xE6, 0x5C, 0x00)
SEV_MEDIUM = RGBColor(0xF9, 0xA8, 0x25)


def new_prs():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)
    return prs


def add_blank_slide(prs):
    layout = prs.slide_layouts[6]  # Blank
    return prs.slides.add_slide(layout)


def set_slide_bg(slide, color):
    bg = slide.background
    fill = bg.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_textbox(slide, left, top, width, height, text, font_size=18,
                bold=False, color=DARK_TEXT, alignment=PP_ALIGN.LEFT,
                font_name="Calibri"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(font_size)
    p.font.bold = bold
    p.font.color.rgb = color
    p.font.name = font_name
    p.alignment = alignment
    return txBox


def add_bullet_list(slide, left, top, width, height, items, font_size=16,
                    color=DARK_TEXT, bold_leads=False, spacing=Pt(6)):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top),
                                     Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True

    for i, item in enumerate(items):
        if i == 0:
            p = tf.paragraphs[0]
        else:
            p = tf.add_paragraph()

        p.space_after = spacing

        if bold_leads and ":" in item:
            # Split on first colon for bold lead
            lead, rest = item.split(":", 1)
            run1 = p.add_run()
            run1.text = lead + ":"
            run1.font.size = Pt(font_size)
            run1.font.bold = True
            run1.font.color.rgb = color
            run1.font.name = "Calibri"
            run2 = p.add_run()
            run2.text = rest
            run2.font.size = Pt(font_size)
            run2.font.bold = False
            run2.font.color.rgb = color
            run2.font.name = "Calibri"
        else:
            run = p.add_run()
            run.text = item
            run.font.size = Pt(font_size)
            run.font.color.rgb = color
            run.font.name = "Calibri"

        p.level = 0
    return txBox


def add_level_bar(slide, top=6.6):
    """Add the 5-level color bar at the bottom of a slide."""
    labels = ["L1: Discovery", "L2: Summarization", "L3: Operational",
              "L4: Investigative", "L5: Adjudicative"]
    bar_width = 2.2
    gap = 0.15
    total = 5 * bar_width + 4 * gap
    start_x = (13.333 - total) / 2

    for i in range(5):
        x = start_x + i * (bar_width + gap)
        shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(x), Inches(top), Inches(bar_width), Inches(0.45)
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = LEVEL_COLORS[i]
        shape.line.fill.background()
        tf = shape.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = labels[i]
        p.font.size = Pt(12)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER


def add_table(slide, left, top, width, col_widths, headers, rows,
              font_size=14):
    """Add a styled table to the slide."""
    num_rows = len(rows) + 1
    num_cols = len(headers)
    table_shape = slide.shapes.add_table(
        num_rows, num_cols, Inches(left), Inches(top),
        Inches(width), Inches(0.45 * num_rows)
    )
    table = table_shape.table

    # Set column widths
    for i, w in enumerate(col_widths):
        table.columns[i].width = Inches(w)

    # Header row
    for i, header in enumerate(headers):
        cell = table.cell(0, i)
        cell.text = header
        cell.fill.solid()
        cell.fill.fore_color.rgb = TABLE_HEADER_BG
        for p in cell.text_frame.paragraphs:
            p.font.size = Pt(font_size)
            p.font.bold = True
            p.font.color.rgb = WHITE
            p.font.name = "Calibri"
            p.alignment = PP_ALIGN.CENTER
        cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    # Data rows
    for r, row in enumerate(rows):
        for c, val in enumerate(row):
            cell = table.cell(r + 1, c)
            cell.text = str(val)
            if r % 2 == 1:
                cell.fill.solid()
                cell.fill.fore_color.rgb = TABLE_ALT_BG
            else:
                cell.fill.solid()
                cell.fill.fore_color.rgb = WHITE
            for p in cell.text_frame.paragraphs:
                p.font.size = Pt(font_size)
                p.font.color.rgb = DARK_TEXT
                p.font.name = "Calibri"
                p.alignment = PP_ALIGN.CENTER
            cell.vertical_anchor = MSO_ANCHOR.MIDDLE

    return table_shape


def add_speaker_notes(slide, text):
    notes_slide = slide.notes_slide
    tf = notes_slide.notes_text_frame
    tf.text = text


def add_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_rounded_rect(slide, left, top, width, height, color):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape


def add_confidential_footer(slide):
    add_textbox(slide, 0.5, 7.05, 3, 0.35, "Confidential  |  March 2026",
                font_size=10, color=MEDIUM_TEXT)


# ── SLIDE BUILDERS ──────────────────────────────────────────────────────

def slide_01_title(prs):
    slide = add_blank_slide(prs)
    # Navy top half
    add_rect(slide, 0, 0, 13.333, 4.5, DARK_NAVY)
    # Title
    add_textbox(slide, 1.5, 0.8, 10.3, 1.5,
                "AI Agent Fidelity\nfor Government",
                font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    # Subtitle line
    add_textbox(slide, 2.5, 2.6, 8.3, 0.6,
                "A Five-Level Framework",
                font_size=28, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)
    # Findings line
    add_textbox(slide, 2, 3.3, 9.3, 0.5,
                "Findings from 462 Test Runs Across 21 Agent Configurations and 6 Models",
                font_size=16, color=RGBColor(0x88, 0xAA, 0xCC),
                alignment=PP_ALIGN.CENTER)
    # Slide count + date
    add_textbox(slide, 4, 4.8, 5.3, 0.5,
                "28 Slides  |  March 2026",
                font_size=16, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)
    # Level bar
    add_level_bar(slide, top=5.6)
    # Stats line
    add_textbox(slide, 2, 6.3, 9.3, 0.4,
                "462 test runs  |  21 agents  |  6 models  |  28 slides  |  2 use cases",
                font_size=14, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)
    add_speaker_notes(slide,
        "Some of you may remember me from such demos as the delegation demo, "
        "or one of my many ALM demos. Those demos and this one have something "
        "in common: I like to build test harnesses for real-world use cases my "
        "customers care about. And they all deal with the same fundamental "
        "topic: accuracy.")


def slide_02_at_a_glance(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "At a Glance",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.3, 10.9, 0.8,
                "462 test runs across 21 agent configurations and 6 models, 2 government use cases, "
                "and 7 testing rounds. Three gaps emerged at the higher levels: tools, "
                "data, and model. Every gap was fixable, and none required AI expertise.",
                font_size=16, color=DARK_TEXT)

    # 2x2 level tiles
    tiles = [
        (L2_LIME, RGBColor(0xF1, 0xF8, 0xE9), "Levels 1-2", "Just Works",
         "8/10 with zero customization.\nModel choice irrelevant."),
        (L3_ORANGE, RGBColor(0xFF, 0xF3, 0xE0), "Level 3", "Data Over Documents",
         "MCP outperformed document search\non every aggregate query."),
        (L4_DEEP_ORANGE, RGBColor(0xFB, 0xE9, 0xE7), "Level 4", "The Inflection Point",
         "Three gaps emerged: tools, data,\nand model. All fixable."),
        (L5_RED, RGBColor(0xFC, 0xE4, 0xEC), "Level 5", "Human in the Loop",
         "High fidelity at L4 means humans\nreview conclusions, not raw data."),
    ]

    for i, (color, tint, label, title, body) in enumerate(tiles):
        col = i % 2
        row = i // 2
        x = 1.2 + col * 4.1
        y = 2.4 + row * 1.65
        add_rounded_rect(slide, x, y, 3.8, 1.4, tint)
        # Color stripe
        add_rect(slide, x, y, 0.15, 1.4, color)
        # Label
        add_textbox(slide, x + 0.3, y + 0.1, 3.3, 0.3, label,
                    font_size=10, bold=True, color=color)
        # Title
        add_textbox(slide, x + 0.3, y + 0.35, 3.3, 0.35, title,
                    font_size=15, bold=True, color=DARK_TEXT)
        # Body
        add_textbox(slide, x + 0.3, y + 0.7, 3.3, 0.6, body,
                    font_size=12, color=MEDIUM_TEXT)

    # Legend on right side
    legend_x = 9.6
    legend_y = 2.4
    add_rounded_rect(slide, legend_x, legend_y, 2.8, 3.3,
                     RGBColor(0xF0, 0xF4, 0xF8))
    add_rect(slide, legend_x, legend_y, 0.1, 3.3, DARK_NAVY)
    add_textbox(slide, legend_x + 0.25, legend_y + 0.15, 2.3, 0.3,
                "The Five Levels", font_size=11, bold=True, color=DARK_NAVY)
    legend_colors = [L1_GREEN, L2_LIME, L3_ORANGE, L4_DEEP_ORANGE, L5_RED]
    legend_names = ["Discovery", "Summarization", "Operational",
                    "Investigative", "Adjudicative"]
    for i, (lc, ln) in enumerate(zip(legend_colors, legend_names)):
        by = legend_y + 0.6 + i * 0.5
        dot = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(legend_x + 0.25), Inches(by), Inches(0.35), Inches(0.35))
        dot.fill.solid()
        dot.fill.fore_color.rgb = lc
        dot.line.fill.background()
        tf = dot.text_frame
        p = tf.paragraphs[0]
        p.text = f"L{i+1}"
        p.font.size = Pt(9)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
        add_textbox(slide, legend_x + 0.7, by, 1.8, 0.35, ln,
                    font_size=12, color=DARK_TEXT)

    # Use case tiles at bottom
    uc_data = [
        (ACCENT_BLUE, "Use Case 1", "Legal Case Analysis",
         "50 cases, 277 people, 333 events", "15 agent configurations"),
        (ACCENT_BLUE, "Use Case 2", "Investigative Analytics",
         "34M rows, 584K properties, 1.6M violations", "8 agent configurations"),
    ]
    for i, (color, label, title, stats, agents) in enumerate(uc_data):
        x = 1.2 + i * 5.6
        y = 5.85
        tint = RGBColor(0xE8, 0xF0, 0xF8)
        add_rounded_rect(slide, x, y, 5.2, 0.95, tint)
        add_rect(slide, x, y, 0.15, 0.95, color)
        add_textbox(slide, x + 0.3, y + 0.05, 2, 0.25, label,
                    font_size=9, bold=True, color=color)
        add_textbox(slide, x + 0.3, y + 0.27, 4.6, 0.3, title,
                    font_size=14, bold=True, color=DARK_TEXT)
        add_textbox(slide, x + 0.3, y + 0.55, 4.6, 0.2, stats,
                    font_size=10, color=MEDIUM_TEXT)
        add_textbox(slide, x + 0.3, y + 0.72, 4.6, 0.2, agents,
                    font_size=10, bold=True, color=DARK_TEXT)

    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "This is the whole story on one slide. Everything that follows is how we got here.")


def slide_03_two_use_cases(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Two Use Cases",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Two real government scenarios, twenty-one agent configurations",
                font_size=22, bold=True, color=DARK_TEXT)

    # Two-column layout
    # UC1
    add_rounded_rect(slide, 1.2, 2.3, 5.4, 3.5, RGBColor(0xE8, 0xF0, 0xF8))
    add_rect(slide, 1.2, 2.3, 0.15, 3.5, ACCENT_BLUE)
    add_textbox(slide, 1.5, 2.4, 4.8, 0.4, "Legal Case Analysis",
                font_size=20, bold=True, color=DARK_NAVY)
    add_bullet_list(slide, 1.5, 3.0, 4.8, 2.5, [
        "50 synthetic legal cases",
        "277 people, 333 timeline events",
        "338 statements, 151 discrepancies",
        "Can an AI agent prepare a case for an "
        "attorney as well as a paralegal?",
    ], font_size=14, color=DARK_TEXT)

    # UC2
    add_rounded_rect(slide, 6.9, 2.3, 5.4, 3.5, RGBColor(0xE8, 0xF0, 0xF8))
    add_rect(slide, 6.9, 2.3, 0.15, 3.5, ACCENT_BLUE)
    add_textbox(slide, 7.2, 2.4, 4.8, 0.4, "Investigative Analytics",
                font_size=20, bold=True, color=DARK_NAVY)
    add_bullet_list(slide, 7.2, 3.0, 4.8, 2.5, [
        "34 million rows of property data",
        "584K properties, 1.6M violations",
        "Can an AI agent surface patterns "
        "an investigator would otherwise miss?",
    ], font_size=14, color=DARK_TEXT)

    add_textbox(slide, 1.2, 6.0, 10.9, 0.4,
                "Both use cases came from real customer conversations. "
                "The data is synthetic but the structure mirrors production systems.",
                font_size=15, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_04_how_built(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "How This Was Built",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "One solution engineer. One AI coding assistant. Four weeks.",
                font_size=22, bold=True, color=DARK_TEXT)

    # Three columns
    cols = [
        (MEDIUM_BLUE, "Human-Directed",
         "Architecture decisions, Azure\nservice selection, schema design,\n"
         "security posture, and the five-level\nfidelity framework came from\n"
         "domain expertise and customer\nconversations."),
        (ACCENT_BLUE, "AI-Accelerated",
         "TypeScript Functions, MCP server,\nSQL schema, 50 synthetic cases,\n"
         "web UI, six PDF generators, and\nthis deck were generated by AI\n"
         "and reviewed by a human."),
        (RGBColor(0x42, 0x95, 0x88), "Human-Verified",
         "462 test runs across 21 agent\nconfigurations and 6 models,\nscored by hand.\n"
         "Every dangerous response was\ncaught through manual review.\n"
         "That is exactly the point\nof Level 5."),
    ]

    for i, (color, title, desc) in enumerate(cols):
        x = 1.2 + i * 3.9
        add_rounded_rect(slide, x, 2.3, 3.5, 3.5, color)
        add_textbox(slide, x + 0.2, 2.5, 3.1, 0.5, title,
                    font_size=20, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, 3.2, 3.1, 2.2, desc,
                    font_size=14, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 6.0, 10.9, 0.4,
                "We practiced what we preach. This entire project is Level 5: "
                "AI-accelerated, human-verified.",
                font_size=15, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_05_meet_agents(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Meet the Agents",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.3, 10.9, 0.6,
                "21 configurations, 3 approaches, 1 set of test prompts",
                font_size=22, bold=True, color=DARK_TEXT)

    # Three approach columns
    col_data = [
        (L1_GREEN, "Zero Code",
         "M365 Copilot",
         "Three JSON manifest files\npointing at the MCP server.\nPlatform picks the model.",
         "1 configuration"),
        (L3_ORANGE, "Configure It",
         "Copilot Studio",
         "SharePoint docs, knowledge base,\nMCP server, or Dataverse --\n"
         "same platform, different\ndata sources.",
         "12 configurations"),
        (ACCENT_BLUE, "Build It Yourself",
         "Pro-Code Agents",
         "Case Analyst (TypeScript)\nInvestigative Agent (OpenAI SDK)\n"
         "Triage Agent (Semantic Kernel)\nFoundry Agent (Agent Service)",
         "8 configurations"),
    ]

    for i, (color, approach, platform, desc, count) in enumerate(col_data):
        x = 1.0 + i * 3.9
        # Card background
        add_rounded_rect(slide, x, 2.1, 3.6, 3.8, color)
        # Approach label
        add_textbox(slide, x + 0.2, 2.2, 3.2, 0.5, approach,
                    font_size=20, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        # Platform
        add_textbox(slide, x + 0.2, 2.75, 3.2, 0.4, platform,
                    font_size=16, color=RGBColor(0xDD, 0xEE, 0xFF),
                    alignment=PP_ALIGN.CENTER)
        # Description
        add_textbox(slide, x + 0.2, 3.25, 3.2, 1.8, desc,
                    font_size=13, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        # Config count pill
        pill_w = 2.4
        pill_x = x + (3.6 - pill_w) / 2
        pill = add_rounded_rect(slide, pill_x, 5.2, pill_w, 0.45,
                                RGBColor(0xFF, 0xFF, 0xFF))
        tf = pill.text_frame
        tf.word_wrap = False
        p = tf.paragraphs[0]
        p.text = count
        p.font.size = Pt(13)
        p.font.bold = True
        p.font.color.rgb = color
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

    add_textbox(slide, 1.2, 6.05, 10.9, 0.4,
                "Same prompts. Same ground truth. "
                "The only variable is how you build it.",
                font_size=15, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "You'll see these names on every chart going forward. The important thing "
        "is the spectrum: on the left, you're up and running in an afternoon. On the "
        "right, you have full control. The question is which level of fidelity "
        "you need -- and that's what the rest of this talk is about.")


def slide_06_the_question(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Question",
                font_size=32, bold=True, color=DARK_NAVY)
    # Accent bar
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    # Big quote
    add_textbox(slide, 1.5, 1.6, 10.3, 1.0,
                '"Not all AI use cases are created equal."',
                font_size=28, bold=True, color=DARK_NAVY,
                alignment=PP_ALIGN.CENTER)

    add_bullet_list(slide, 1.5, 2.8, 10.3, 2.5, [
        "Your policy lookup chatbot and your legal case preparation tool have "
        "fundamentally different accuracy requirements",
        "This framework helps you decide where to deploy aggressively, where to "
        "add guardrails, and where human review is non-negotiable",
    ], font_size=18, color=DARK_TEXT)

    # Arrow visual: Low Stakes -> High Stakes
    add_rect(slide, 1.5, 5.0, 10.3, 0.08, RGBColor(0xDD, 0xDD, 0xDD))
    # Gradient dots
    for i, c in enumerate(LEVEL_COLORS):
        x = 1.5 + i * 2.58
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(x + 0.95), Inches(4.75), Inches(0.5), Inches(0.5))
        dot.fill.solid()
        dot.fill.fore_color.rgb = c
        dot.line.fill.background()
    add_textbox(slide, 1.5, 5.2, 3, 0.4, "Low Stakes",
                font_size=14, color=L1_GREEN, alignment=PP_ALIGN.LEFT)
    add_textbox(slide, 8.8, 5.2, 3, 0.4, "High Stakes",
                font_size=14, color=L5_RED, alignment=PP_ALIGN.RIGHT)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "A policy lookup chatbot that hallucinates a deadline is annoying. A legal case "
        "prep tool that misses a discrepancy can change an outcome. The framework we're "
        "about to walk through helps you match your engineering investment to the stakes "
        "of the use case -- so you know when 'good enough' really is good enough, and "
        "when it isn't.")


def slide_03_five_levels(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Five Levels",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    headers = ["Level", "Name", "Stakes", "Example"]
    rows = [
        ["1", "Discovery", "Inconvenience", '"Find our leave policy"'],
        ["2", "Summarization", "Wasted time", '"Summarize this audit report"'],
        ["3", "Operational", "Misallocated resources", '"How many open cases by type?"'],
        ["4", "Investigative", "Missed evidence", '"Build a timeline from these records"'],
        ["5", "Adjudicative", "Wrong legal outcome", '"Prepare facts for this hearing"'],
    ]
    tbl = add_table(slide, 1.2, 1.6, 10.9, [1.0, 2.2, 3.0, 4.7], headers, rows, font_size=15)

    # Color the level number cells
    table = tbl.table
    for r in range(5):
        cell = table.cell(r + 1, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = LEVEL_COLORS[r]
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True

    add_textbox(slide, 1.2, 5.3, 10.9, 0.5,
                'Where does your highest-priority use case fall?',
                font_size=18, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_04_quick_win(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Levels 1-2: The Quick Win",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.7,
                "Document agents score 8 out of 10 with zero engineering",
                font_size=24, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.4, 10.9, 3.0, [
        "Point Copilot at your SharePoint library",
        "No custom code, no databases, no tool design",
        "Model choice barely matters at this level",
        "Good enough for policy lookup, report summarization, meeting notes",
    ], font_size=18)

    # Simple flow diagram: Copilot -> SharePoint -> Answer
    boxes = [("Copilot", 3.0), ("SharePoint", 5.8), ("Answer", 8.6)]
    for label, x in boxes:
        shape = add_rounded_rect(slide, x, 4.6, 2.0, 0.8, ACCENT_BLUE)
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = label
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER
    # Arrows between boxes
    for x in [5.05, 7.85]:
        add_textbox(slide, x, 4.65, 0.7, 0.8, ">",
                    font_size=28, bold=True, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 5.8, 10.9, 0.5,
                "This is where most agencies should start. It works today.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_05_level3(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Level 3: Structured Data Starts Winning",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.7,
                "Aggregate queries work across all agent types",
                font_size=24, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.4, 10.9, 1.8, [
        "Document agents excel at single-case analysis (10/10) but can't count, "
        "sum, or filter across a full dataset",
        '"How many active cases?" "What\'s the breakdown by county?" "Top 5 violators?" '
        "-- these require structured data",
        "Agents with structured data access -- both Copilot Studio MCP and pro-code -- "
        "answer every aggregate query correctly when paired with a capable model",
    ], font_size=18)

    # Side-by-side comparison boxes
    # Doc agent box
    box_l = add_rounded_rect(slide, 1.5, 4.2, 4.8, 1.8, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, 1.7, 4.25, 4.4, 0.4, "Document Agent", font_size=14,
                bold=True, color=L5_RED)
    add_textbox(slide, 1.7, 4.7, 4.4, 1.0,
                '"I don\'t have that information."',
                font_size=16, color=MEDIUM_TEXT, alignment=PP_ALIGN.CENTER)

    # MCP agent box
    box_r = add_rounded_rect(slide, 7.0, 4.2, 4.8, 1.8, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, 7.2, 4.25, 4.4, 0.4, "MCP Agent", font_size=14,
                bold=True, color=L1_GREEN)
    # Mini table text
    add_textbox(slide, 7.2, 4.7, 4.4, 1.0,
                "CPS: 18 cases\nTPR: 12 cases\nGuardianship: 8 cases",
                font_size=14, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 6.2, 10.9, 0.4,
                "This is where you start investing in data structure.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_06_inflection(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Level 4: The Inflection Point",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.7,
                "This is where engineering decisions determine success or failure",
                font_size=22, bold=True, color=DARK_TEXT)

    add_textbox(slide, 1.2, 2.3, 10.9, 0.5,
                "Three gaps emerged in testing:",
                font_size=18, color=DARK_TEXT)

    # Three gap cards
    gaps = [
        ("The Model Gap", "GPT-4.1 scores 9.5/10\nvs GPT-4o at 4/10", L4_DEEP_ORANGE),
        ("The Tool Gap", "One missing tool caused\n87% failure rate", L3_ORANGE),
        ("The Data Gap", "Facts buried in narrative\ntext were invisible", MEDIUM_BLUE),
    ]
    for i, (title, desc, color) in enumerate(gaps):
        x = 1.5 + i * 3.8
        box = add_rounded_rect(slide, x, 3.0, 3.3, 2.5, color)
        add_textbox(slide, x + 0.2, 3.1, 2.9, 0.6, title,
                    font_size=20, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.2, 3.8, 2.9, 1.5, desc,
                    font_size=16, color=WHITE,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 5.8, 10.9, 0.4,
                "Level 4 is where the investment pays off -- or doesn't.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_07_model_gap(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Model Gap",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "GPT-4.1 vs GPT-4o: Same tools, same data, same backend",
                font_size=20, bold=True, color=DARK_TEXT)

    # Bar chart visual (horizontal bars)
    models = [
        ("GPT-4.1", 9.5, L1_GREEN),
        ("GPT-4o (Government Cloud)", 4.0, L3_ORANGE),
        ("Platform-assigned (M365)", 2.0, L5_RED),
    ]
    for i, (label, score, color) in enumerate(models):
        y = 2.3 + i * 1.2
        # Label
        add_textbox(slide, 1.5, y, 4.5, 0.4, label, font_size=16,
                    color=DARK_TEXT, alignment=PP_ALIGN.RIGHT)
        # Bar
        bar_max_w = 5.5
        bar_w = bar_max_w * (score / 10.0)
        add_rounded_rect(slide, 6.2, y + 0.05, bar_w, 0.45, color)
        # Score label
        add_textbox(slide, 6.2 + bar_w + 0.1, y, 1.5, 0.5,
                    f"{score}/10", font_size=18, bold=True, color=color)

    add_bullet_list(slide, 1.5, 5.2, 10.3, 1.5, [
        "Government Cloud currently defaults to GPT-4o; expanded model availability on the roadmap",
        "No amount of prompt engineering closed this gap",
        "Pro-code agents can use GPT-4.1 in Government Cloud via Azure OpenAI directly",
    ], font_size=16)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "The model is the single biggest variable at Level 4.")


def slide_08_tool_gap(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Tool Gap",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.7,
                "87% failure rate from one missing tool",
                font_size=24, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.4, 10.9, 2.0, [
        'Agents could not convert "4763 Griscom St" to a database parcel ID',
        "15 address lookups across 5 agents, only 2 succeeded",
        "One fuzzy-match tool added: zero failures in retesting",
        "Investigative Agent: 1 out of 10 to a perfect 10 out of 10",
    ], font_size=18)

    # Before/After visual
    # Before box
    add_rounded_rect(slide, 2.0, 4.6, 4.2, 1.6, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, 2.2, 4.65, 3.8, 0.4, "BEFORE", font_size=14,
                bold=True, color=L5_RED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 2.2, 5.1, 3.8, 0.9,
                "13 of 15 lookups failed\n1/10 accuracy",
                font_size=18, color=L5_RED, alignment=PP_ALIGN.CENTER)

    # Arrow
    add_textbox(slide, 6.2, 4.9, 0.9, 0.8, ">",
                font_size=36, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    # After box
    add_rounded_rect(slide, 7.1, 4.6, 4.2, 1.6, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, 7.3, 4.65, 3.8, 0.4, "AFTER", font_size=14,
                bold=True, color=L1_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.3, 5.1, 3.8, 0.9,
                "0 failures\n10/10 accuracy",
                font_size=18, color=L1_GREEN, alignment=PP_ALIGN.CENTER)

    # Highlight the change
    add_textbox(slide, 2.0, 6.3, 9.3, 0.4,
                "The change: one search_properties fuzzy-match tool",
                font_size=14, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "Tool design is as important as model selection. Maybe more.")


def slide_09_level5(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Level 5: Trust But Verify",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Even our best agent reproduced a dangerous error",
                font_size=22, bold=True, color=DARK_TEXT)

    # Split screen: Sheriff Report vs Medical Records
    # Left box
    add_rounded_rect(slide, 1.2, 2.3, 5.2, 2.0, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, 1.4, 2.35, 4.8, 0.35, "Sheriff's Report",
                font_size=14, bold=True, color=L5_RED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.4, 2.8, 4.8, 1.2,
                '"no fractures detected\non skeletal survey"',
                font_size=18, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    # Right box
    add_rounded_rect(slide, 6.9, 2.3, 5.2, 2.0, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, 7.1, 2.35, 4.8, 0.35, "Medical Records",
                font_size=14, bold=True, color=L1_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.1, 2.8, 4.8, 1.2,
                "Bilateral long bone\nfractures documented",
                font_size=18, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 4.45, 10.9, 0.4,
                "7 of 8 document agents quoted the sheriff's report without cross-checking",
                font_size=16, bold=True, color=SEV_CRITICAL,
                alignment=PP_ALIGN.CENTER)

    # Document improvement section
    add_textbox(slide, 1.2, 5.0, 10.9, 0.4,
                "The fix: document structure improvements",
                font_size=18, bold=True, color=DARK_NAVY)
    add_bullet_list(slide, 1.2, 5.4, 10.9, 1.5, [
        "Added cross-reference headers to each document (zero content changes)",
        "Commercial agent: 8/10 to 10/10 -- pulled Medical Records as primary source",
        "GCC agent: 3/10 to 9/10 across all 10 prompts with document hygiene alone",
        "Zero code, zero engineering -- document hygiene any paralegal can implement",
    ], font_size=14)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "Connect to daily experience: Copilot helps you draft an email -- you review it. "
        "Copilot helps you write code -- you review it. Why would we skip human review at "
        "any of these five levels, especially Level 5 where the stakes are a family's future?")


def slide_10_danger_taxonomy(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Danger Taxonomy",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Five categories of AI failure, ranked by severity",
                font_size=20, bold=True, color=DARK_TEXT)

    dangers = [
        ("1", "CRITICAL", "False Negative",
         "Data retrieved, answer not recognized", SEV_CRITICAL),
        ("2", "CRITICAL", "Faithfully Reproduced Misinformation",
         "Source document had an error, agent quoted it accurately", SEV_CRITICAL),
        ("3", "HIGH", "Misattribution",
         "Real fact, wrong person", SEV_HIGH),
        ("4", "HIGH", "Hallucinated Fact with Confidence",
         "Invented detail with no source", SEV_HIGH),
        ("5", "MEDIUM", "Silent Failure",
         "No results, no indication anything went wrong", SEV_MEDIUM),
    ]

    for i, (num, severity, name, desc, color) in enumerate(dangers):
        y = 2.2 + i * 0.9
        # Severity pill
        pill = add_rounded_rect(slide, 1.5, y + 0.05, 1.5, 0.5, color)
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = severity
        p.font.size = Pt(11)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        # Name + desc
        add_textbox(slide, 3.2, y, 3.5, 0.6, f"{num}. {name}",
                    font_size=16, bold=True, color=DARK_TEXT)
        add_textbox(slide, 6.8, y + 0.05, 5.5, 0.6, desc,
                    font_size=15, color=MEDIUM_TEXT)

    add_textbox(slide, 1.2, 6.0, 10.9, 0.4,
                "These are not hypothetical. We documented each one across 462 test runs.",
                font_size=16, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_11_iterative_process(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Iterative Process",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Deploying an agent is not a one-time event",
                font_size=22, bold=True, color=DARK_TEXT)

    # Four round boxes in a cycle
    rounds = [
        ("Round 0", "Baseline Testing", "Measure failures\nagainst ground truth",
         RGBColor(0x78, 0x90, 0x9C)),
        ("Round 1", "Fix the Data", "Make facts discrete\nand queryable",
         MEDIUM_BLUE),
        ("Round 2", "Fix the Tools", "Help the model\nreach the data",
         ACCENT_BLUE),
        ("Round 3", "Validate Models", "Confirm what works\nwhere",
         RGBColor(0x42, 0x95, 0x88)),
    ]

    for i, (rnd, title, desc, color) in enumerate(rounds):
        x = 1.2 + i * 3.1
        box = add_rounded_rect(slide, x, 2.5, 2.7, 3.0, color)
        add_textbox(slide, x + 0.1, 2.55, 2.5, 0.45, rnd,
                    font_size=14, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 3.1, 2.5, 0.5, title,
                    font_size=18, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 3.7, 2.5, 1.2, desc,
                    font_size=14, color=RGBColor(0xDD, 0xEE, 0xFF),
                    alignment=PP_ALIGN.CENTER)

    # Arrows between
    for i in range(3):
        x = 1.2 + (i + 1) * 3.1 - 0.4
        add_textbox(slide, x, 3.3, 0.5, 0.6, ">",
                    font_size=28, bold=True, color=ACCENT_BLUE,
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 5.8, 10.9, 0.4,
                "Every organization will go through these rounds, in this order.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_12_results(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Results After Iteration",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Six agents reached 9-10 out of 10 after iterative improvement",
                font_size=20, bold=True, color=DARK_TEXT)

    headers = ["Agent", "Round 0", "Final", "Fix"]
    rows = [
        ["SP/PDF/GCC (Copilot Studio)", "3/10", "9/10", "Cross-referenced docs"],
        ["KB/DOCX/Com (Copilot Studio)", "6/10", "10/10", "Cross-referenced docs"],
        ["Commercial MCP (Copilot Studio)", "8/10", "10/10", "Descriptions + data"],
        ["Investigative Agent (OpenAI SDK)", "1/10", "10/10", "Fuzzy-match tool"],
        ["Foundry Agent", "4/10", "9/10", "Descriptions + data"],
        ["Triage Agent (Semantic Kernel)", "0/10", "10/10", "5 rounds: tools+data+model"],
    ]
    add_table(slide, 1.5, 2.3, 10.3, [4.2, 1.5, 1.5, 3.1], headers, rows,
              font_size=14)

    add_textbox(slide, 1.2, 5.5, 10.9, 0.7,
                "The first two rows required zero engineering -- just document hygiene.\n"
                "The rest required iterative tool and data work.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_13_code_spectrum(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Code Spectrum",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Zero code to full code -- same fidelity when you control the model",
                font_size=20, bold=True, color=DARK_TEXT)

    headers = ["Approach", "Code", "Best For"]
    rows = [
        ["M365 Copilot (Com)", "Zero (3 JSON manifests)", "Levels 1-2"],
        ["Copilot Studio", "Zero to low code", "Levels 1-3"],
        ["Foundry Agent", "Minimal code", "Levels 3-4"],
        ["Custom SDK (OpenAI, SK)", "Full code", "Levels 4-5"],
    ]
    add_table(slide, 1.5, 2.3, 10.3, [3.5, 3.5, 3.0], headers, rows,
              font_size=16)

    add_textbox(slide, 1.2, 5.0, 10.9, 0.7,
                "The investment buys governance and customization, not fidelity.\n"
                "Fidelity comes from the model and the data.",
                font_size=18, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_14_why_copilot_studio(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Why Copilot Studio",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "One platform, declarative to pro-code",
                font_size=22, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.3, 10.9, 3.5, [
        "Same environment hosts zero-code SharePoint agents AND MCP-connected structured data agents",
        "Model flexibility: swap GPT-4o for GPT-4.1, accuracy goes from 4/10 to 10/10 -- no config changes",
        "Enterprise governance built in: DLP, audit logging, admin pre-approval for tool calls",
        "M365 distribution: agents show up in Teams and Copilot chat -- no separate app to deploy",
        "Our test data proves the platform works -- the model is the variable, not the architecture",
    ], font_size=17)

    add_textbox(slide, 1.2, 5.5, 10.9, 0.7,
                "Commercial Copilot Studio MCP scored a perfect 10 out of 10.\n"
                "The architecture is sound. The limiting factor is model availability in GCC.",
                font_size=16, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "This slide is critical for the tech series. MUST articulate the platform value, "
        "not just the results.")


def slide_15_what_to_do(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "What to Do at Each Level",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    headers = ["Level", "Action"]
    rows = [
        ["1-2", "Deploy Copilot with SharePoint. Clean up document metadata. Done."],
        ["3", "Connect to structured data via MCP. Clear tool descriptions. Summary modes."],
        ["4", "Purpose-built tools. Models with strong multi-step reasoning (GPT-4.1 in our testing). Ground truth test suite. Budget 3+ rounds."],
        ["5", "Level 4 + human review workflows + audit logging + citation linking."],
    ]
    tbl = add_table(slide, 1.2, 1.6, 10.9, [1.5, 9.0], headers, rows, font_size=16)

    # Color the level cells
    table = tbl.table
    level_colors_map = [L1_GREEN, L3_ORANGE, L4_DEEP_ORANGE, L5_RED]
    for r in range(4):
        cell = table.cell(r + 1, 0)
        cell.fill.solid()
        cell.fill.fore_color.rgb = level_colors_map[r]
        for p in cell.text_frame.paragraphs:
            p.font.color.rgb = WHITE
            p.font.bold = True

    add_textbox(slide, 1.2, 5.2, 10.9, 0.5,
                "Match your investment to your accuracy requirement.",
                font_size=20, bold=True, color=ACCENT_BLUE,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_16_gcc(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Government Cloud Customers",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Copilot Studio works at every level -- the levers differ",
                font_size=22, bold=True, color=DARK_TEXT)

    # Three-tier recommendation boxes
    tiers = [
        (L2_LIME, RGBColor(0xF1, 0xF8, 0xE9),
         "Levels 1-3", "Copilot Studio,\nany configuration"),
        (L4_DEEP_ORANGE, RGBColor(0xFB, 0xE9, 0xE7),
         "Level 4 (Documents)", "Copilot Studio +\ngood document hygiene\n(9/10 proven, GPT-4o, GCC)"),
        (MEDIUM_BLUE, RGBColor(0xE8, 0xF0, 0xF8),
         "Level 4-5 (Structured Data)", "Model selection matters\nSonnet 4.6 = 11/11 (Commercial)\nGPT-4.1 via Azure OpenAI (GCC)"),
    ]

    for i, (color, tint, label, body) in enumerate(tiers):
        x = 1.2 + i * 3.9
        add_rounded_rect(slide, x, 2.3, 3.5, 2.2, tint)
        add_rect(slide, x, 2.3, 0.15, 2.2, color)
        add_textbox(slide, x + 0.3, 2.4, 3.0, 0.4, label,
                    font_size=14, bold=True, color=color)
        add_textbox(slide, x + 0.3, 2.85, 3.0, 1.5, body,
                    font_size=14, color=DARK_TEXT)

    # GCC gap callout
    gap_box = add_rounded_rect(slide, 1.2, 4.8, 10.9, 1.5,
                               RGBColor(0xF5, 0xF5, 0xF5))
    add_rect(slide, 1.2, 4.8, 0.15, 1.5, ACCENT_BLUE)
    add_textbox(slide, 1.6, 4.85, 3.0, 0.35, "The GCC Gap",
                font_size=16, bold=True, color=DARK_NAVY)
    add_bullet_list(slide, 1.6, 5.2, 9.8, 1.0, [
        "Model selection not yet available in GCC Copilot Studio -- document quality is your lever today",
        "When expanded model availability arrives, every agent improves overnight with zero changes",
        "Start building ground truth test suites now so you are ready to measure the improvement",
    ], font_size=13, color=DARK_TEXT, spacing=Pt(4))

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "Copilot Studio is not just for Levels 1-3. We proved 9/10 at Level 4 "
        "with document agents in GCC using GPT-4o -- no code, just better documents. "
        "The gap is only for structured data queries where model selection matters. "
        "Know what works today. Plan for what's coming.")


def slide_16b_scoring(prs):
    """S21: How We Score -- Ground Truth Grading."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "How We Score: Ground Truth Grading",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.35, 10.9, 0.6,
                "Every grade is verifiable -- we wrote the data",
                font_size=22, bold=True, color=DARK_TEXT)

    # ── Left side: 0-10 scale table ──
    scale_headers = ["Score", "Meaning", "Example"]
    scale_rows = [
        ("10", "Fully correct, complete, sourced", "All facts match, sources cited"),
        ("8-9", "Correct, one minor detail missing", "All statements found, missed one page ref"),
        ("5", "Partially correct or incomplete", "Found one of two sources, or wrong time"),
        ("1-4", "Significant gaps", "Right case, missed key facts"),
        ("0", "Wrong or dangerously confident", "Real citation, wrong conclusion"),
    ]
    scale_col_widths = [0.7, 2.5, 2.6]
    add_table(slide, 1.2, 2.2, 5.8, scale_col_widths,
              scale_headers, scale_rows, font_size=12)

    # ── Right side: Process steps ──
    add_textbox(slide, 7.5, 2.2, 5.0, 0.4,
                "The Process",
                font_size=18, bold=True, color=DARK_NAVY)
    add_rect(slide, 7.5, 2.6, 4.5, 0.04, ACCENT_BLUE)

    process_items = [
        "Each prompt scored 0-10 against ground truth",
        "One run per agent, per-prompt scores averaged",
        "Same prompt to every agent -- side by side",
        "Graded by Claude against ground truth (live)",
        "462 test runs, 21 agents, 6 models",
    ]
    add_bullet_list(slide, 7.5, 2.8, 5.0, 2.8, process_items,
                    font_size=14, color=DARK_TEXT, spacing=Pt(8))

    # ── Bottom callout box ──
    callout = add_rounded_rect(slide, 1.2, 5.4, 10.9, 1.0,
                               RGBColor(0xE8, 0xF0, 0xF8))
    add_rect(slide, 1.2, 5.4, 0.15, 1.0, ACCENT_BLUE)
    add_textbox(slide, 1.6, 5.5, 10.2, 0.8,
                '"We grade against data we wrote. Every correct answer is '
                'verifiable. Every wrong answer is provably wrong -- not a '
                'matter of interpretation."',
                font_size=16, color=MEDIUM_TEXT)

    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "All 50 cases are synthetic -- we control the ground truth for every prompt. "
        "The 0-10 scale distinguishes between an agent that found the right case but "
        "missed one detail (8-9) versus one that got half the answer (5) versus one "
        "that confidently gave the wrong answer (0). During the live demo, I paste "
        "agent responses into Claude Code and it grades them in real time.")


def slide_17_demo(prs):
    slide = add_blank_slide(prs)
    # Full navy background
    set_slide_bg(slide, DARK_NAVY)

    add_textbox(slide, 1.5, 0.6, 10.3, 1.0,
                "Live Demo",
                font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.5, 1.6, 10.3, 0.6,
                '"Same agent. Change one variable. Watch the result change."',
                font_size=22, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)

    # Part 1: Document Quality (left side)
    part1_x = 1.0
    add_textbox(slide, part1_x, 2.5, 5.5, 0.5,
                "Part 1: Document Quality  (~10 min)",
                font_size=18, bold=True, color=WHITE)
    add_rect(slide, part1_x, 3.0, 5.5, 0.04, ACCENT_BLUE)

    part1_items = [
        "3 agents, same model, 3 SharePoint libraries",
        "Raw  ->  Cross-Referenced  ->  Enriched",
        "Skeletal survey: FAIL -> PASS -> PASS",
        "Marcus Webb: FAIL -> PARTIAL -> PASS",
        "Score arc: 3/10 -> 7/10 -> 8/10",
        "Cost: $0, 45 minutes of document hygiene",
    ]
    for i, item in enumerate(part1_items):
        add_textbox(slide, part1_x + 0.3, 3.2 + i * 0.55, 5.2, 0.5,
                    item, font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Part 2: Model Selection (right side)
    part2_x = 7.0
    add_textbox(slide, part2_x, 2.5, 5.5, 0.5,
                "Part 2: Model Selection  (~8 min)",
                font_size=18, bold=True, color=WHITE)
    add_rect(slide, part2_x, 3.0, 5.5, 0.04, ACCENT_BLUE)

    part2_items = [
        "1 Dataverse MCP agent, swap model live",
        "GPT-4o  ->  Sonnet 4.6",
        "TPR case filter: 0 -> 10",
        "Cross-doc reasoning: 0 -> 10",
        "Time gap calculation: 0 -> 10",
        "Full battery: 3.2/11 -> 6/11 -> 11/11",
    ]
    for i, item in enumerate(part2_items):
        add_textbox(slide, part2_x + 0.3, 3.2 + i * 0.55, 5.2, 0.5,
                    item, font_size=14, color=RGBColor(0xCC, 0xDD, 0xEE))

    # Bottom tagline
    add_textbox(slide, 1.5, 6.5, 10.3, 0.5,
                "All demos run in Copilot Studio -- no custom code UI",
                font_size=14, bold=True, color=RGBColor(0x88, 0xAA, 0xCC),
                alignment=PP_ALIGN.CENTER)


def slide_17b_one_prompt(prs):
    """S22: One Prompt, Three Interventions — Marcus Webb progression."""
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "One Prompt, Three Interventions",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    # The prompt
    prompt_box = add_rounded_rect(slide, 1.2, 1.35, 10.9, 0.7,
                                   RGBColor(0xF5, 0xF5, 0xF5))
    add_textbox(slide, 1.5, 1.4, 10.3, 0.6,
                '"What did Marcus Webb tell hospital staff about when he put Jaylen '
                'to bed, and did he give the same answer to law enforcement?"',
                font_size=15, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    # Three columns: Raw, Cross-Ref, Enriched
    col_configs = [
        (SEV_CRITICAL, RGBColor(0xFF, 0xEB, 0xEE), "Raw Documents",
         "FAIL",
         "Found: Sheriff Report interview",
         "Missed: Nursing notes in Medical Records",
         "The agent has no reason to look\nin a medical record for\nparent statements."),
        (SEV_MEDIUM, RGBColor(0xFF, 0xF8, 0xE1), "Cross-Referenced",
         "PARTIAL",
         "Found: Both sources identified",
         "Missed: Specific nursing note detail",
         "Cross-reference header points to\nMedical Records, but the agent\nstill can't find nursing notes."),
        (L1_GREEN, RGBColor(0xE8, 0xF5, 0xE9), "Enriched + Metadata",
         "PASS",
         "Found: Both statements, correct times",
         "Missed: Nothing",
         'SP metadata: "nursing interview,\nparent statements" on Medical\nRecords enables retrieval.'),
    ]

    for i, (color, tint, title, result, found, missed, explanation) in enumerate(col_configs):
        x = 0.9 + i * 4.1
        # Card background
        add_rounded_rect(slide, x, 2.25, 3.7, 3.7, tint)
        # Color stripe at top
        add_rect(slide, x, 2.25, 3.7, 0.12, color)
        # Title
        add_textbox(slide, x + 0.15, 2.45, 3.4, 0.35, title,
                    font_size=16, bold=True, color=DARK_TEXT,
                    alignment=PP_ALIGN.CENTER)
        # Result pill
        pill = add_rounded_rect(slide, x + 1.15, 2.85, 1.4, 0.4, color)
        tf = pill.text_frame
        p = tf.paragraphs[0]
        p.text = result
        p.font.size = Pt(14)
        p.font.bold = True
        p.font.color.rgb = WHITE
        p.font.name = "Calibri"
        p.alignment = PP_ALIGN.CENTER

        # Found / Missed
        add_textbox(slide, x + 0.2, 3.4, 3.3, 0.45, found,
                    font_size=12, color=DARK_TEXT)
        add_textbox(slide, x + 0.2, 3.85, 3.3, 0.45, missed,
                    font_size=12, bold=True, color=color)
        # Explanation
        add_textbox(slide, x + 0.2, 4.4, 3.3, 1.2, explanation,
                    font_size=11, color=MEDIUM_TEXT)

    # Arrows between columns
    for i in range(2):
        ax = 4.55 + i * 4.1
        arrow = slide.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW,
            Inches(ax), Inches(3.7), Inches(0.5), Inches(0.35)
        )
        arrow.fill.solid()
        arrow.fill.fore_color.rgb = ACCENT_BLUE
        arrow.line.fill.background()

    # Bottom callout
    callout = add_rounded_rect(slide, 1.5, 6.15, 10.3, 0.55,
                                RGBColor(0xE3, 0xF2, 0xFD))
    add_rect(slide, 1.5, 6.15, 0.12, 0.55, ACCENT_BLUE)
    add_textbox(slide, 1.8, 6.2, 9.8, 0.45,
                "Cross-references fix reasoning errors. "
                "Metadata fixes retrieval errors. You need both.",
                font_size=16, bold=True, color=DARK_NAVY,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "This is the most relatable prompt in the deck. Everyone understands "
        "'did the parent tell the same story to the nurse and the police?' "
        "Raw docs: agent only found the Sheriff Report. Cross-refs: found both "
        "sources but couldn't retrieve the nursing notes detail. Metadata: the "
        "keywords 'nursing interview' and 'parent statements' on Medical Records "
        "told the agent exactly where to look. $0, 45 minutes, no code.")


def slide_18_challenged_premise(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Surprising Finding: Agent Challenged Its Own Premise",
                font_size=28, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "A pro-code agent questioned the question -- and was right",
                font_size=20, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.3, 10.9, 2.5, [
        "Asked about a $146,000 purchase vs a $357,100 assessment",
        "Agent checked the source data: actual assessment was $53,155",
        "Recalculated using verified numbers instead of the prompt's numbers",
        "No other agent questioned the input",
    ], font_size=18)

    # Quote box
    quote_box = add_rounded_rect(slide, 2.0, 4.6, 9.3, 1.2,
                                  RGBColor(0xE3, 0xF2, 0xFD))
    add_textbox(slide, 2.2, 4.65, 8.9, 1.1,
                '"The assessment value in our records is $53,155, not $357,100 as stated '
                'in the question. Using the verified assessment value..."',
                font_size=15, color=MEDIUM_BLUE, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 6.0, 10.9, 0.4,
                "This is what custom orchestration enables -- but only with iterative testing.",
                font_size=16, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_19_false_negative(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Surprising Finding: False Negative",
                font_size=28, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "The model retrieved the answer and didn't recognize it",
                font_size=20, bold=True, color=DARK_TEXT)

    add_bullet_list(slide, 1.2, 2.3, 10.9, 2.0, [
        "Agent called the right tools",
        'Received data containing "two positive drug screens (fentanyl)"',
        'Concluded: "no drug test results exist in the available data"',
        "Invisible to users: tool calls look correct, answer sounds confident",
    ], font_size=18)

    # Tool output vs Agent conclusion
    add_rounded_rect(slide, 1.5, 4.5, 5.0, 1.5, RGBColor(0xE8, 0xF5, 0xE9))
    add_textbox(slide, 1.7, 4.55, 4.6, 0.35, "Tool Output (correct)",
                font_size=13, bold=True, color=L1_GREEN, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 1.7, 5.0, 4.6, 0.8,
                '"two positive drug screens\n(fentanyl)"',
                font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_rounded_rect(slide, 6.8, 4.5, 5.0, 1.5, RGBColor(0xFF, 0xEB, 0xEE))
    add_textbox(slide, 7.0, 4.55, 4.6, 0.35, "Agent Conclusion (wrong)",
                font_size=13, bold=True, color=L5_RED, alignment=PP_ALIGN.CENTER)
    add_textbox(slide, 7.0, 5.0, 4.6, 0.8,
                '"no drug test results exist\nin the available data"',
                font_size=16, color=DARK_TEXT, alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 6.2, 10.9, 0.4,
                "This is why you need ground truth testing. You cannot catch this in production.",
                font_size=16, bold=True, color=SEV_CRITICAL,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)


def slide_20_bottom_line(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "The Bottom Line",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.5, 10.9, 0.5,
                "Three numbers to remember",
                font_size=22, bold=True, color=DARK_TEXT)

    # Three big number cards
    cards = [
        ("8/10", "What you get with\nzero engineering", "Levels 1-2",
         L1_GREEN),
        ("9.5/10", "What you get with\npurpose-built tools + GPT-4.1", "Levels 3-4",
         ACCENT_BLUE),
        ("Trust\nBut Verify", "No score is high enough\nto skip human review", "Level 5",
         L5_RED),
    ]

    for i, (number, desc, levels, color) in enumerate(cards):
        x = 1.2 + i * 3.9
        box = add_rounded_rect(slide, x, 2.3, 3.5, 3.5, color)
        add_textbox(slide, x + 0.1, 2.5, 3.3, 1.2, number,
                    font_size=40, bold=True, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 3.8, 3.3, 1.0, desc,
                    font_size=15, color=WHITE,
                    alignment=PP_ALIGN.CENTER)
        add_textbox(slide, x + 0.1, 5.0, 3.3, 0.5, levels,
                    font_size=13, color=RGBColor(0xDD, 0xEE, 0xFF),
                    alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.2, 6.0, 10.9, 0.5,
                '"The question is not \'should we deploy AI?\' The question is '
                '\'which level are we operating at, and have we invested accordingly?\'"',
                font_size=15, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_level_bar(slide)
    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "So what do we do? Do we just give up on agents? No -- you continue to test "
        "them and make them better, and you always adhere to trust but verify. I have "
        "agents building apps for me right now. If I were building production enterprise "
        "apps, I'd have a team of human developers checking the work, helping test, "
        "reviewing agentic results. The AI accelerates the team. The team validates the AI.")


def slide_24_scorecard(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, WHITE)
    add_textbox(slide, 0.8, 0.4, 11.7, 0.7,
                "Full Scorecard",
                font_size=32, bold=True, color=DARK_NAVY)
    add_rect(slide, 0.8, 1.1, 2, 0.06, ACCENT_BLUE)

    add_textbox(slide, 1.2, 1.3, 10.9, 0.5,
                "462 test runs, 21 agent configurations, 6 models, 2 use cases",
                font_size=20, bold=True, color=DARK_TEXT)

    # UC1 table
    add_textbox(slide, 1.2, 1.9, 5.5, 0.35,
                "Use Case 1: Legal Case Analysis (16 agents)",
                font_size=13, bold=True, color=DARK_NAVY)
    uc1_headers = ["Agent", "Model", "Score"]
    uc1_widths = [2.6, 1.3, 0.8]
    uc1_rows = [
        ["Case Analyst (pro-code)", "GPT-4.1", "10/10"],
        ["CS MCP/Com", "GPT-4.1", "10/10"],
        ["CS SP/PDF/Com", "GPT-4.1", "10/10"],
        ["CS KB/DOCX/Com", "GPT-4.1", "10/10"],
        ["CS DV/Com", "Sonnet 4.6", "10/10"],
        ["CS DV/Com", "GPT-5 Reasoning", "10/10"],
        ["CS SP/DOCX/Com", "GPT-4.1", "9/10"],
        ["CS MCP/GCC", "GPT-4o", "9/10"],
        ["CS SP/PDF/GCC", "GPT-4o", "9/10"],
        ["CS KB/PDF/GCC", "GPT-4o", "9/10"],
        ["CS KB/PDF/Com", "GPT-4.1", "8/10"],
        ["CS KB/DOCX/GCC", "GPT-4o", "8/10"],
        ["CS SP/DOCX/GCC", "GPT-4o", "7/10"],
        ["CS DV/Com", "GPT-4.1", "6/10"],
        ["CS DV/Com", "GPT-5 Auto", "4/10"],
        ["CS DV/GCC", "GPT-4o", "3.2/10"],
    ]
    add_table(slide, 0.8, 2.25, 4.7, uc1_widths, uc1_headers, uc1_rows,
              font_size=9)

    # UC2 table
    add_textbox(slide, 7.0, 1.9, 5.5, 0.35,
                "Use Case 2: Investigative Analytics (8 agents)",
                font_size=13, bold=True, color=DARK_NAVY)
    uc2_headers = ["Agent", "Model", "Score"]
    uc2_widths = [2.6, 1.3, 0.8]
    uc2_rows = [
        ["CS MCP/Com", "GPT-4.1", "10/10"],
        ["Investigative (OpenAI SDK)", "GPT-4.1", "10/10"],
        ["Triage (Semantic Kernel)", "GPT-4.1", "10/10"],
        ["Foundry Agent", "GPT-4.1", "9/10"],
        ["CS SP/PDF/GCC", "GPT-4o", "9/10"],
        ["CS SP/PDF/Com", "GPT-4.1", "8/10"],
        ["CS MCP/GCC", "GPT-4o", "4/10"],
        ["M365 Copilot MCP/Com", "Platform", "2/10"],
    ]
    add_table(slide, 6.6, 2.25, 4.7, uc2_widths, uc2_headers, uc2_rows,
              font_size=9)

    add_textbox(slide, 1.2, 6.6, 10.9, 0.4,
                "Dataverse MCP (5 models): Sonnet 4.6 (10/10), GPT-5 Reasoning (10/10), GPT-4.1 (6/10), "
                "GPT-5 Auto (4/10), GPT-4o (3.2/10). Reasoning models close the gap.",
                font_size=12, bold=True, color=MEDIUM_TEXT,
                alignment=PP_ALIGN.CENTER)

    add_confidential_footer(slide)
    add_speaker_notes(slide,
        "This is 462 data points. The pattern is clear: "
        "the model and the data matter more than the platform.")


def slide_25_next_steps(prs):
    slide = add_blank_slide(prs)
    set_slide_bg(slide, DARK_NAVY)

    add_textbox(slide, 1.5, 0.8, 10.3, 1.0,
                "Next Steps",
                font_size=44, bold=True, color=WHITE,
                alignment=PP_ALIGN.CENTER)

    add_textbox(slide, 1.5, 2.0, 10.3, 0.6,
                "Start the conversation",
                font_size=24, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)

    questions = [
        "Where do your use cases fall on the spectrum?",
        "What data do you already have in structured form?",
        "What's your risk tolerance for AI-generated output?",
        "Ready for a deeper dive? We have the architecture, the test data, and the framework.",
    ]

    for i, q in enumerate(questions):
        y = 3.2 + i * 0.8
        add_textbox(slide, 3.0, y, 7.3, 0.6, q,
                    font_size=18, color=WHITE)
        # bullet dot
        dot = slide.shapes.add_shape(
            MSO_SHAPE.OVAL, Inches(2.5), Inches(y + 0.15), Inches(0.2), Inches(0.2))
        dot.fill.solid()
        dot.fill.fore_color.rgb = ACCENT_BLUE
        dot.line.fill.background()

    add_textbox(slide, 1.5, 6.2, 10.3, 0.6,
                "I'm not here to sell you a product. I'm here to help you make "
                "the right investment for your mission.",
                font_size=16, bold=True, color=RGBColor(0xA0, 0xC4, 0xE8),
                alignment=PP_ALIGN.CENTER)


def main():
    prs = new_prs()

    slide_01_title(prs)              # S1  Title
    slide_02_at_a_glance(prs)        # S2  At a Glance
    slide_03_two_use_cases(prs)      # S3  The Two Use Cases
    slide_04_how_built(prs)          # S4  How This Was Built
    slide_05_meet_agents(prs)        # S5  Meet the Agents
    slide_06_the_question(prs)       # S6  The Question
    slide_03_five_levels(prs)        # S7  The Five Levels
    slide_04_quick_win(prs)          # S8  Levels 1-2
    slide_05_level3(prs)             # S9  Level 3
    slide_06_inflection(prs)         # S10 Level 4
    slide_07_model_gap(prs)          # S11 Model Gap
    slide_08_tool_gap(prs)           # S12 Tool Gap
    slide_09_level5(prs)             # S13 Level 5
    slide_10_danger_taxonomy(prs)    # S14 Danger Taxonomy
    slide_11_iterative_process(prs)  # S15 Iterative Process
    slide_12_results(prs)            # S16 Results After Iteration
    slide_13_code_spectrum(prs)      # S17 Code Spectrum
    slide_14_why_copilot_studio(prs) # S18 Why Copilot Studio
    slide_15_what_to_do(prs)         # S19 What to Do
    slide_16_gcc(prs)                # S20 GCC
    slide_16b_scoring(prs)           # S21 How We Score
    slide_17_demo(prs)               # S22 Live Demo
    slide_17b_one_prompt(prs)        # S23 One Prompt, Three Interventions
    slide_18_challenged_premise(prs) # S24 Surprising: Premise
    slide_19_false_negative(prs)     # S25 Surprising: False Negative
    slide_20_bottom_line(prs)        # S26 The Bottom Line
    slide_24_scorecard(prs)          # S27 Full Scorecard
    slide_25_next_steps(prs)         # S28 Next Steps

    out_dir = os.path.join(os.path.dirname(__file__), "..", "decks")
    os.makedirs(out_dir, exist_ok=True)
    out_path = os.path.join(out_dir, "agent-fidelity-spectrum.pptx")
    prs.save(out_path)
    print(f"Saved: {out_path}")
    print(f"Slides: {len(prs.slides)}")


if __name__ == "__main__":
    main()
